import os
import openai
import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import pytesseract
from PIL import Image
import time
import PyPDF2
import io


st.set_page_config(
    page_title="Summarizer",
    page_icon="📒",
    layout="wide",
)

session_state = st.session_state

if "api_key" not in session_state:
    session_state.api_key = None

# Check if the provided OpenAI API key is valid
def is_valid_openai_key(api_key):
    """Check if the provided OpenAI API key is valid."""
    openai.api_key = api_key
    try:
        # Make a test call.
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert in summarizing Documents."},
            ],
        )
        if response and 'choices' in response:
            return True
    except openai.error.AuthenticationError:
        return False
    return False

# If API key has not been set in the code, show the input
if not session_state.api_key:
    st.title("📒Welcome to AI summarizer website.")
    st.write("Please enter your OpenAI API key to begin:")
    temp_key = st.text_input("OpenAI API Key:", type="password")
    api_key_url = "https://www.maisieai.com/help/how-to-get-an-openai-api-key-for-chatgpt"
    html_code = f"""
    <ul style="list-style-type: none; padding-left: 0;">
    <li>Need an API key? <a href="{api_key_url}" target="_blank">Learn how to obtain one here.</a></li>
    </ul>
    """
    st.markdown(html_code, unsafe_allow_html=True)

    if st.button("Submit"):
        if len(temp_key) < 20:
            st.error("The API key seems too short. Please recheck.")
        else:
            if is_valid_openai_key(temp_key):
                session_state.api_key = temp_key
                st.experimental_rerun()
            else:
                st.error("The provided API key is invalid. Please recheck.")
openai.api_key = session_state.api_key
if session_state.api_key:

    def get_response(text): 
         # Check if the input exceeds the maximum allowed context length
        max_context_length = 4000
        if len(text.split()) > max_context_length:
            st.error(f"Input exceeds the maximum allowed length of {max_context_length} words. Please provide a shorter input.")
            return None  
        prompt = f"""
            You are an expert in summarizing Documents. You will be given a Document delimited by four backquotes, 
            make sure to capture the main points, key arguments, and any supporting evidence presented in the article.
            Your summary should be informative and well-structured, ideally consisting of 3-5 sentences.

            text: ''''{text}''''
            """
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": prompt},
                ],
            )
            return response
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
            return None       

else:
    st.error("Please provide a valid API key.")


def extract_text_from_docx(file):
    # Check if the file is a Word document
    if not file.name.lower().endswith(".docx"):
        st.error("Invalid file type. Please upload a Word document.")
        return None

    doc = docx.Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)



def extract_text_from_JPG(file):
    # Check if the file is an image
    if not file.type.startswith("image"):
        st.error("Invalid file type. Please upload an image.")
        return None

    image = Image.open(file)
    pytesseract.pytesseract.tesseract_cmd = r'C:\Users\mishal\miniconda3\envs\py311\Lib\site-packages\tesseract.exe'
    raw_text = pytesseract.image_to_string(image)
    return raw_text  



def main():

    st.divider()
    st.markdown(
        "*Second step*: Choose format and insert a file to summarize it.")
    option = st.radio("Select Input Type",
                      ("Text", "Image", "PDF", "Word", "PowerPoint"))

    if option == "Text":
        if not session_state.api_key:
            st.error("Please enter an API key first.")
            return

        user_input = st.text_area(
            "Enter Text", placeholder="Enter some paragraphs to summarize it."
        )
        if st.button("Submit", key=12) and user_input != "":
            max_context_length = 4000
            if len(user_input.split()) > max_context_length:
                st.error(f"Input exceeds the maximum allowed length of {max_context_length} words. Please provide a shorter input.")
                return
            else:
                 progress_text = "In progress. Please wait."
                 my_bar = st.progress(0, text=progress_text)

            for percent_complete in range(100):
                time.sleep(0.1)
                my_bar.progress(percent_complete + 1, text=progress_text)

            time.sleep(1)
            my_bar.empty()

            response = get_response(user_input)
            generated_response = response['choices'][0]['message']['content']
            generated_text = generated_response.strip()
            st.success("done!")
            st.subheader("Summary")
            st.markdown(f"> {generated_text}")

        else:
            st.error("Please enter some text.")
        
    elif option in ["Image", "PDF", "Word", "PowerPoint"]:
        if not session_state.api_key:
            st.error("Please enter an API key first.")
            return

        uploaded_file = st.file_uploader(f"Choose a {option} file", type=["jpg", "pdf", "docx", "pptx"])

        st.divider()

        if st.button("Submit") and uploaded_file is not None:
            text = extract_content_from_file(option, uploaded_file)

            if text is None:
                return
            
            max_context_length = 4000
            if len(text.split()) > max_context_length:
                st.error(f"Input exceeds the maximum allowed length of {max_context_length} words. Please provide a shorter input.")
                return
            else:
                progress_text = "In progress. Please wait."
                my_bar = st.progress(0, text=progress_text)

            for percent_complete in range(100):
                time.sleep(0.1)
                my_bar.progress(percent_complete + 1, text=progress_text)

            time.sleep(1)
            my_bar.empty()

            response = get_response(text=text)
            summary = response['choices'][0]['message']['content']
            st.success("done!")
            st.subheader("Summary")
            st.markdown(f"> {summary}")

       
    
    st.divider()


    st.subheader("IMPORTANT NOTE:")
    st.caption(""" This is a demo version, and more features are on the horizon! In the upcoming semester we will do the following:

- *Enhanced Functionality:* Additional features to make your summarization experience even better.
- *User-Friendly Design:* A sleek and intuitive design for a more pleasant interaction.
- *Expanded Format Support:* More file types and document formats for greater flexibility.
- *Optimized Performance:* Faster processing and improved efficiency.
- *Community Feedback Integration:* Your feedback will shape the next versions!

Stay tuned for the next release and the exciting improvements. Thank you for using our AI Summarizer! """)
    st.caption('v1.0')

def extract_content_from_file(option, file):
    if option == "Image":
        return extract_text_from_JPG(file)
    elif option == "PDF":
        return extract_text_from_pdf(file)
    elif option == "Word":
        return extract_text_from_docx(file)
    elif option == "PowerPoint":
        return extract_text_from_pptx(file)
    else:
        
        return None
    

def extract_text_from_pdf(file):
    # Check if the file is a PDF
    if not file.name.lower().endswith(".pdf"):
        st.error("Invalid file type. Please upload a PDF file.")
        return None

    
    file_content = io.BytesIO(file.read())

    try:
        pdf_reader = PyPDF2.PdfReader(file_content)
        text = ""
        for page_number in range(len(pdf_reader.pages)):
            text += pdf_reader.pages[page_number].extract_text()
        return text
    except PyPDF2.utils.PdfReadError:
        st.error("Error reading PDF file.")
        return None



def extract_text_from_pptx(file):
    # Check if the file is a PowerPoint 
    if not file.name.lower().endswith(".pptx"):
        st.error("Invalid file type. Please upload a PowerPoint.")
        return None

    raw_text = ""
    prs = Presentation(file)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                raw_text = raw_text + text

    return raw_text


if __name__ == "__main__":
    main()